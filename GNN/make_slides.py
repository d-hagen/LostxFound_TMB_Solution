"""Generate a skeleton .pptx for the 5-year TMB lost-and-found plan."""
from pptx import Presentation
from pathlib import Path

OUT = Path(__file__).parent / "TMB_5yr_plan.pptx"

prs = Presentation()
title_layout = prs.slide_layouts[0]
content_layout = prs.slide_layouts[1]


def title_slide(title, subtitle=""):
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title
    if subtitle:
        s.placeholders[1].text = subtitle


def bullet_slide(title, bullets):
    s = prs.slides.add_slide(content_layout)
    s.shapes.title.text = title
    tf = s.placeholders[1].text_frame
    tf.text = bullets[0]
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = b


# ── 1. Title ─────────────────────────────────────────────────────────────
title_slide(
    "TMB Lost & Found — 5-Year Plan",
    "From single-station depot to model-routed locker network",
)

# ── 2. Current situation ─────────────────────────────────────────────────
bullet_slide(
    "Current situation",
    [
        "All items centralised at Sagrada Família depot",
        "Registration delay: 5–10 days after find",
        "Claims only in person, by phone, or via web form",
        "No data captured on passenger preference or convenience",
    ],
)

# ── 3. Pre-GNN changes ───────────────────────────────────────────────────
bullet_slide(
    "Proposed changes — before the model",
    [
        "Unified database: TMB + City Council + Mossos + Renfe",
        "Staff mobile app: photo → VLM auto-description → instant registration",
        "Chatbot retrieval: natural-language search against the shared ledger",
        "Pickup questionnaire: capture preferred station + item type",
        "Outcome: searchable items in seconds, easier retrieval, labelled data",
    ],
)

# ── 4. GNN application ───────────────────────────────────────────────────
bullet_slide(
    "GNN application",
    [
        "Input:  found-at station, lost-at station (claimant-reported), item type, hour, day-of-week",
        "Output: probability distribution over the 157 stations — the claimant's preferred pickup location",
        "Decision rule: choose the storage station that minimises expected transfer-aware travel (E[hops]) from the predicted pickup distribution",
        "Aggregating chosen storage stations across the event corpus sizes the locker network — how many lockers, and where to place them",
    ],
)

# ── 5. How it's trained ──────────────────────────────────────────────────
bullet_slide(
    "How the GNN is trained",
    [
        "Year 0: collect questionnaire labels at existing pickup points",
        "Pretrain on prior synth data (encoded passenger-profile assumptions)",
        "Fine-tune on real questionnaire data; prior weight decays year by year",
        "Annual retraining on accumulated questionnaire corpus",
    ],
)

# ── 6. What it tells / how it changes the system ─────────────────────────
bullet_slide(
    "What it tells us and what changes",
    [
        "Per-item: ship directly to the locker the model recommends",
        "Aggregate: which stations host lockers, sized by predicted weekly load",
        "Replaces 'everything to Sagrada' with model-driven routing at find time",
        "Locker map reshapes annually as travel patterns drift",
    ],
)

# ── 7. Hopeful impact ────────────────────────────────────────────────────
bullet_slide(
    "Hopeful impact",
    [
        "Convenient pickup → more trust → higher retrieval rate",
        "Faster retrieval → less storage time per item",
        "Positive feedback loop: better data → better model → better placement",
    ],
)

# ── 8. TMB tradeoffs ─────────────────────────────────────────────────────
bullet_slide(
    "TMB — tradeoffs",
    [
        "+  Shorter average storage time per item",
        "+  Fewer items held at any one moment",
        "+  Higher passenger trust in the service",
        "−  More distributed storage logistics",
        "−  Operational complexity of locker network upkeep",
    ],
)

# ── 9. Users ─────────────────────────────────────────────────────────────
bullet_slide(
    "Users",
    [
        "+  Closer pickup points",
        "+  Anytime pickup via locker code (no info-point queue, no opening hours)",
        "+  Less anxiety about the process",
    ],
)

# ── 10. Demo idea ────────────────────────────────────────────────────────
bullet_slide(
    "Demo idea — full live system",
    [
        "Two laptops, end-to-end flow with audience participation",
        "",
        "Laptop 1 — Item registration (VLM)",
        "    Scan any item the audience hands over",
        "    VLM generates description; entry added to the shared database",
        "    Random found-at station/line assigned for the demo",
        "    Barcelona map shows where the GNN will store the item",
        "",
        "Laptop 2 — Retrieval (LLM chatbot)",
        "    Passenger describes the lost item in natural language",
        "    LLM matches against the database and returns the pickup locker",
        "    Same Barcelona map highlights the storage station",
    ],
)

prs.save(OUT)
print(f"Saved {OUT}  ({prs.slide_width}x{prs.slide_height}, {len(prs.slides)} slides)")
